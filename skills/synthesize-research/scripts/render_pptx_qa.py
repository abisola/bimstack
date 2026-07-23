#!/usr/bin/env python3
"""Faithful-ish layout renderer: reads the saved .pptx geometry/fills/text/images
and rasterises each slide with PIL so layout issues (overlap, overflow, off-slide)
can be visually QA'd without LibreOffice."""
import os, io, glob
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

HERE = os.path.dirname(__file__)
import sys; PPTX = sys.argv[1]
OUT = sys.argv[2]; os.makedirs(OUT, exist_ok=True)
for f in glob.glob(os.path.join(OUT, "*.png")): os.remove(f)

PXI = 132  # px per inch
EMU_IN = 914400
def px(emu): return int(round((emu / EMU_IN) * PXI))

GEO   = "/System/Library/Fonts/Supplemental/Georgia.ttf"
GEO_B = "/System/Library/Fonts/Supplemental/Georgia Bold.ttf"
GEO_I = "/System/Library/Fonts/Supplemental/Georgia Italic.ttf"
ARI   = "/System/Library/Fonts/Supplemental/Arial.ttf"
ARI_B = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"
ARI_I = "/System/Library/Fonts/Supplemental/Arial Italic.ttf"
_cache = {}
def font(path, size):
    k=(path,size)
    if k not in _cache: _cache[k]=ImageFont.truetype(path, size)
    return _cache[k]
FIG_R=os.path.join(HERE,"fonts","Figtree-Regular.ttf")
FIG_S=os.path.join(HERE,"fonts","Figtree-SemiBold.ttf")
def pick(fontname, bold, ital, size):
    fn = fontname or ""
    if "Figtree" in fn and os.path.exists(FIG_R):
        p = FIG_S if bold else FIG_R
    elif "Georgia" in fn:
        p = GEO_B if bold else (GEO_I if ital else GEO)
        if ital and not os.path.exists(GEO_I): p = GEO
    else:
        p = ARI_B if bold else (ARI_I if ital else ARI)
        if ital and not os.path.exists(ARI_I): p = ARI
    return font(p, max(8, int(round(size*PXI/72.0))))

def rgb(c):
    try: return (c[0], c[1], c[2])
    except Exception: return None

def shape_fill(sp):
    try:
        f = sp.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            return rgb(f.fore_color.rgb)
    except Exception:
        return None
    return None

prs = Presentation(PPTX)
SW = px(prs.slide_width); SH = px(prs.slide_height)

def wrap_runs(draw, runs, fnt_for, maxw):
    """Return list of lines; each line is list of (text, font, color)."""
    lines=[]; cur=[]; curw=0
    space_cache={}
    for text, fnt, col in runs:
        # split keeping explicit breaks
        for bi, segment in enumerate(text.split("\n")):
            if bi>0:
                lines.append(cur); cur=[]; curw=0
            for word in segment.split(" "):
                if word=="":
                    w_word=draw.textlength(" ", font=fnt);
                    cur.append((" ",fnt,col)); curw+=w_word; continue
                ww = draw.textlength(word, font=fnt)
                sw = draw.textlength(" ", font=fnt)
                add = ww + (sw if cur else 0)
                if cur and curw+add > maxw:
                    lines.append(cur); cur=[(word,fnt,col)]; curw=ww
                else:
                    if cur: cur.append((" ",fnt,col)); curw+=sw
                    cur.append((word,fnt,col)); curw+=ww
    if cur: lines.append(cur)
    return lines

for idx, slide in enumerate(prs.slides, 1):
    img = Image.new("RGB",(SW,SH),(255,255,255))
    d = ImageDraw.Draw(img)
    for sp in slide.shapes:
        x,y,w,h = px(sp.left or 0), px(sp.top or 0), px(sp.width or 0), px(sp.height or 0)
        # picture
        if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                im = Image.open(io.BytesIO(sp.image.blob)).convert("RGBA")
                im = im.resize((max(1,w),max(1,h)))
                img.paste(im,(x,y),im)
            except Exception as e:
                d.rectangle([x,y,x+w,y+h], outline=(200,200,200))
            continue
        # table
        if getattr(sp, "has_table", False) and sp.has_table:
            tb = sp.table
            colw = [px(c.width) for c in tb.columns]; roww = [px(r.height) for r in tb.rows]
            cy0 = y
            for ri, row in enumerate(tb.rows):
                cx0 = x
                for ci, cell in enumerate(row.cells):
                    cw_=colw[ci]; ch_=roww[ri]
                    cf=None
                    try:
                        if cell.fill.type is not None and cell.fill.fore_color and cell.fill.fore_color.type is not None:
                            cf=rgb(cell.fill.fore_color.rgb)
                    except Exception: cf=None
                    if cf is not None: d.rectangle([cx0,cy0,cx0+cw_,cy0+ch_], fill=cf)
                    d.rectangle([cx0,cy0,cx0+cw_,cy0+ch_], outline=(216,210,210), width=1)
                    para=cell.text_frame.paragraphs[0]
                    runs=para.runs
                    if runs:
                        r0=runs[0]; sz=r0.font.size.pt if r0.font.size else 9
                        try: col=rgb(r0.font.color.rgb) if (r0.font.color and r0.font.color.type is not None) else (40,40,40)
                        except Exception: col=(40,40,40)
                        f=pick(r0.font.name,bool(r0.font.bold),bool(r0.font.italic),sz)
                        t=cell.text; tw=d.textlength(t,font=f)
                        al=str(para.alignment or "")
                        if "CENTER" in al: tx2=cx0+(cw_-tw)//2
                        elif "RIGHT" in al: tx2=cx0+cw_-tw-px(0.05)
                        else: tx2=cx0+px(0.07)
                        d.text((tx2, cy0+(ch_- (sz*PXI/72))//2), t, font=f, fill=col)
                    cx0+=cw_
                cy0+=roww[ri]
            continue
        # fill
        fill = shape_fill(sp)
        try: is_oval = "OVAL" in str(sp.auto_shape_type or "")
        except Exception: is_oval = False
        if fill is not None:
            if is_oval: d.ellipse([x,y,x+w,y+h], fill=fill)
            else: d.rectangle([x,y,x+w,y+h], fill=fill)
        # line/outline
        try:
            if sp.line and sp.line.color and sp.line.color.type is not None:
                lc = rgb(sp.line.color.rgb)
                if lc and fill is None:
                    d.rectangle([x,y,x+w,y+h], outline=lc, width=2)
        except Exception: pass
        # text
        if sp.has_text_frame and sp.text_frame.text.strip():
            tf = sp.text_frame
            ml=px(tf.margin_left or 0); mr=px(tf.margin_right or 0)
            mt=px(tf.margin_top or 0); mb=px(tf.margin_bottom or 0)
            tx=x+ml; tw=w-ml-mr
            # gather paragraphs
            para_lines=[]  # (lines, align, space_after_px, approx_line_h)
            for para in tf.paragraphs:
                runs=[]
                align = str(para.alignment) if para.alignment else "LEFT"
                maxsz=12
                for r in para.runs:
                    sz = r.font.size.pt if r.font.size else 18
                    maxsz=max(maxsz,sz)
                    try:
                        col = rgb(r.font.color.rgb) if (r.font.color and r.font.color.type is not None) else (40,40,40)
                    except Exception:
                        col = (40,40,40)
                    fnt = pick(r.font.name, bool(r.font.bold), bool(r.font.italic), sz)
                    runs.append((r.text, fnt, col or (40,40,40)))
                if not runs:
                    para_lines.append(([], align, 6, int(maxsz*PXI/72*1.2))); continue
                lines = wrap_runs(d, runs, None, tw)
                lh = int(maxsz*PXI/72*1.18)
                sa = px(para.space_after) if para.space_after else int(0.04*PXI)
                ls = para.line_spacing if para.line_spacing else 1.0
                if isinstance(ls,(int,float)) and ls>0: lh=int(lh*ls)
                para_lines.append((lines, align, sa, lh))
            # vertical anchor
            total_h=sum(len(pl[0])*pl[3]+pl[2] for pl in para_lines)
            va = tf.vertical_anchor
            if va is not None and str(va)=="MIDDLE (3)" or (va is not None and "MIDDLE" in str(va)):
                cy = y + max(0,(h-total_h)//2)
            elif va is not None and "BOTTOM" in str(va):
                cy = y + max(0,(h-total_h))
            else:
                cy = y
            overflow = total_h > h+4
            for lines, align, sa, lh in para_lines:
                for line in lines:
                    lw=sum(d.textlength(t,font=f) for t,f,c in line)
                    if align=="CENTER (2)" or align=="CENTER": lx=tx+max(0,(tw-lw)//2)
                    elif align=="RIGHT (3)" or align=="RIGHT": lx=tx+max(0,tw-lw)
                    else: lx=tx
                    cx=lx
                    for t,f,c in line:
                        d.text((cx,cy), t, font=f, fill=c)
                        cx+=d.textlength(t,font=f)
                    cy+=lh
                cy+=sa
            if overflow:
                d.rectangle([x,y,x+w,y+h], outline=(255,0,255), width=3)
    img.save(os.path.join(OUT,f"slide-{idx:02d}.png"))
    print(f"slide-{idx:02d}.png  ({SW}x{SH})")
print("DONE")
